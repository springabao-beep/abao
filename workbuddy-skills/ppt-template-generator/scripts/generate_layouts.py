# -*- coding: utf-8 -*-
"""
Career Compass PPT Generator
Generates PPT slides following the Career Compass Visual Guidelines.

Usage:
    from generate_layouts import (
        CareerCompassPPTX,
        cover_slide, three_col_points, pillars_showcase,
        timeline, data_showcase, detailed_content,
        discussion_slide, team_grid, quote_slide, closing_slide
    )

    pptx = CareerCompassPPTX()
    cover_slide(pptx, "CAREER COMPASS", "2026 Q3 Talent Review")
    pptx.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Career Compass Color System ─────────────────────────────────────────────
CC_PRIMARY   = RGBColor(0x9D, 0x3D, 0x1D)   # #9D3D1D — titles, emphasis, dark bg
CC_SECONDARY = RGBColor(0xB9, 0x5B, 0x42)   # #B95B42 — subtitles
CC_ACCENT    = RGBColor(0xD6, 0x75, 0x5A)   # #D6755A — body emphasis
CC_LIGHT     = RGBColor(0xDD, 0xBC, 0xB0)   # #DDBCB0 — card backgrounds
CC_BASE      = RGBColor(0xF1, 0xE2, 0xDD)   # #F1E2DD — large-area light fills
CC_PAGE      = RGBColor(0xF7, 0xF3, 0xED)   # #F7F3ED — page base
CC_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CC_BLACK     = RGBColor(0x00, 0x00, 0x00)

# ─── Typography ────────────────────────────────────────────────────────────────
FONT_BODY   = "Futura Bk BT"   # primary font
FONT_FALLBACK = "Arial"          # system fallback
FONT_CN     = "Microsoft YaHei"  # Chinese fallback

# Slide dimensions (16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)


def new_presentation():
    """Create a new 16:9 Career Compass branded Presentation."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _add_slide(prs):
    """Add a slide using the first available layout (blank-friendly)."""
    try:
        # Try index 0 first (usually Title Slide or Blank)
        return prs.slides.add_slide(prs.slide_layouts[0])
    except (IndexError, TypeError):
        return prs.slides.add_slide(None)


# ─── Utility Helpers ─────────────────────────────────────────────────────────

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=CC_BLACK,
                 align=PP_ALIGN.LEFT, font_name=FONT_BODY):
    """Add a text box to a slide with Career Compass defaults."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line=False):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line:
        shape.line.color.rgb = CC_BLACK
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def set_slide_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def slide_number(slide, number, dark_bg=False):
    """Add slide number bottom-right."""
    color = CC_WHITE if dark_bg else CC_BLACK
    add_text_box(slide, str(number), 12.5, 7.0, 0.7, 0.35,
                 font_size=11, color=color, align=PP_ALIGN.RIGHT)


# ─── Layout 1: Cover Slide ────────────────────────────────────────────────────
def cover_slide(pptx, title, subtitle="", number=None):
    """
    Dark cover: #9D3D1D background, centered 54pt white bold title,
    28pt light subtitle. Optional slide number bottom-right.
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_PRIMARY)

    add_text_box(slide, title, 1.0, 2.5, 11.33, 1.5,
                 font_size=54, bold=True, color=CC_WHITE,
                 align=PP_ALIGN.CENTER, font_name=FONT_BODY)

    if subtitle:
        add_text_box(slide, subtitle, 1.5, 4.2, 10.33, 0.8,
                     font_size=28, bold=False, color=CC_LIGHT,
                     align=PP_ALIGN.CENTER, font_name=FONT_BODY)

    if number:
        slide_number(slide, number, dark_bg=True)

    return slide


# ─── Layout 2: Three-Column Key Points ──────────────────────────────────────
def three_col_points(pptx, heading, points, number=None):
    """
    Heading + 3 equal cards (CC_BASE background), each with:
    - Number badge (CC_PRIMARY circle)
    - Bold title (CC_PRIMARY)
    - Description text (black)
    points: list of dicts [{title, description}, ...] — max 3
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_PAGE)

    add_text_box(slide, heading, 0.5, 0.3, 12.33, 0.8,
                 font_size=32, bold=True, color=CC_BLACK)

    card_top = 1.4
    card_h   = 5.0
    card_w   = 3.8
    card_gap = 0.3
    start_x  = 0.5

    for i, pt in enumerate(points[:3]):
        x = start_x + i * (card_w + card_gap)

        # Card background
        add_rect(slide, x, card_top, card_w, card_h, CC_BASE)

        # Number circle
        circle = slide.shapes.add_shape(
            9,  # MSO_SHAPE_TYPE.OVAL
            Inches(x + 0.1), Inches(card_top + 0.1),
            Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = CC_PRIMARY
        circle.line.fill.background()

        add_text_box(slide, str(i + 1),
                     x + 0.1, card_top + 0.12, 0.5, 0.46,
                     font_size=20, bold=True, color=CC_WHITE,
                     align=PP_ALIGN.CENTER)

        # Point title
        add_text_box(slide, pt.get("title", ""), x + 0.2, card_top + 0.75,
                     card_w - 0.4, 0.6,
                     font_size=20, bold=True, color=CC_BLACK)

        # Description
        add_text_box(slide, pt.get("description", ""), x + 0.2,
                     card_top + 1.45, card_w - 0.4, 3.2,
                     font_size=14, color=CC_BLACK)

    if number:
        slide_number(slide, number)

    return slide


# ─── Layout 3: Pillars Showcase ─────────────────────────────────────────────
def pillars_showcase(pptx, title, pillars, number=None):
    """
    Heading + 2×2 grid of pillar cards (CC_LIGHT bg) with left accent bar (CC_PRIMARY).
    pillars: list of dicts [{title, description}, ...] — max 4
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_PAGE)

    add_text_box(slide, title, 0.5, 0.3, 12.33, 0.8,
                 font_size=32, bold=True, color=CC_BLACK)

    card_w  = 5.5
    card_h  = 2.4
    gap_x   = 0.8
    gap_y   = 0.3
    start_x = 0.5
    start_y = 1.3

    positions = [
        (start_x, start_y),
        (start_x + card_w + gap_x, start_y),
        (start_x, start_y + card_h + gap_y),
        (start_x + card_w + gap_x, start_y + card_h + gap_y),
    ]

    for i, (x, y) in enumerate(positions[:4]):
        pt = pillars[i] if i < len(pillars) else {}
        add_rect(slide, x, y, card_w, card_h, CC_LIGHT)
        # Left accent bar
        add_rect(slide, x, y, 0.08, card_h, CC_PRIMARY)
        # Title
        add_text_box(slide, pt.get("title", ""), x + 0.2, y + 0.2,
                     card_w - 0.3, 0.5,
                     font_size=18, bold=True, color=CC_BLACK)
        # Description
        add_text_box(slide, pt.get("description", ""), x + 0.2, y + 0.8,
                     card_w - 0.3, 1.4,
                     font_size=13, color=CC_BLACK)

    if number:
        slide_number(slide, number)

    return slide


# ─── Layout 4: Timeline / Journey ───────────────────────────────────────────
def timeline(pptx, title, steps, number=None):
    """
    Heading + horizontal timeline with n connected circular nodes.
    Each node: number inside CC_PRIMARY circle + title + description below.
    steps: list of dicts [{number, title, description}, ...] — max 5
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_PAGE)

    add_text_box(slide, title, 0.5, 0.3, 12.33, 0.7,
                 font_size=32, bold=True, color=CC_BLACK)

    n       = min(len(steps), 5)
    node_w  = 2.0
    node_gap = (12.33 - n * node_w) / (n + 1)
    node_top = 2.8
    node_sz  = 0.7

    for i, step in enumerate(steps[:n]):
        x = node_gap + i * (node_w + node_gap)

        # Circle node
        circle = slide.shapes.add_shape(
            9,  # oval
            Inches(x + (node_w - node_sz) / 2),
            Inches(node_top),
            Inches(node_sz),
            Inches(node_sz))
        circle.fill.solid()
        circle.fill.fore_color.rgb = CC_PRIMARY
        circle.line.fill.background()

        # Number inside
        add_text_box(slide, str(step.get("number", i + 1)),
                     x + (node_w - node_sz) / 2, node_top + 0.1,
                     node_sz, node_sz,
                     font_size=18, bold=True, color=CC_WHITE,
                     align=PP_ALIGN.CENTER)

        # Connector line (not last)
        if i < n - 1:
            line = slide.shapes.add_shape(
                1,  # rectangle = line
                Inches(x + node_w),
                Inches(node_top + node_sz / 2 - 0.02),
                Inches(node_gap + 0.08),
                Inches(0.04))
            line.fill.solid()
            line.fill.fore_color.rgb = CC_SECONDARY
            line.line.fill.background()

        # Step title
        add_text_box(slide, step.get("title", ""), x, node_top + 0.9,
                     node_w, 0.5,
                     font_size=13, bold=True, color=CC_BLACK,
                     align=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, step.get("description", ""), x, node_top + 1.4,
                     node_w, 2.0,
                     font_size=11, color=CC_BLACK,
                     align=PP_ALIGN.CENTER)

    if number:
        slide_number(slide, number)

    return slide


# ─── Layout 5: Data Showcase ────────────────────────────────────────────────
def data_showcase(pptx, title, metrics, number=None):
    """
    Heading + up to 3 metric cards (CC_BASE bg), each with:
    - Large value (72pt bold CC_BLACK)
    - Label (18pt bold black)
    - Optional sublabel (13pt black)
    metrics: list of dicts [{value, label, sublabel?}, ...] — max 3
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_PAGE)

    add_text_box(slide, title, 0.5, 0.3, 12.33, 0.7,
                 font_size=32, bold=True, color=CC_BLACK)

    n       = min(len(metrics), 3)
    card_w  = 3.8
    card_gap = 0.4
    start_x = 0.5
    card_top = 1.5
    card_h  = 5.0

    for i, m in enumerate(metrics[:n]):
        x = start_x + i * (card_w + card_gap)
        add_rect(slide, x, card_top, card_w, card_h, CC_BASE)

        add_text_box(slide, m.get("value", ""), x, card_top + 0.5,
                     card_w, 1.5,
                     font_size=72, bold=True, color=CC_BLACK,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, m.get("label", ""), x + 0.2, card_top + 2.2,
                     card_w - 0.4, 1.0,
                     font_size=18, bold=True, color=CC_BLACK,
                     align=PP_ALIGN.CENTER)

        if m.get("sublabel"):
            add_text_box(slide, m.get("sublabel", ""), x + 0.2, card_top + 3.3,
                         card_w - 0.4, 0.8,
                         font_size=13, color=CC_BLACK,
                         align=PP_ALIGN.CENTER)

    if number:
        slide_number(slide, number)

    return slide


# ─── Layout 6: Detailed Content ─────────────────────────────────────────────
def detailed_content(pptx, title, body_text, info_box_text="", number=None):
    """
    Two-column: body text area (left, CC_BASE) + info/summary box (right, CC_LIGHT).
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_PAGE)

    add_text_box(slide, title, 0.5, 0.3, 12.33, 0.7,
                 font_size=32, bold=True, color=CC_BLACK)

    add_rect(slide, 0.5, 1.2, 7.5, 5.5, CC_BASE)
    add_text_box(slide, body_text, 0.7, 1.4, 7.1, 5.1,
                 font_size=16, color=CC_BLACK)

    if info_box_text:
        add_rect(slide, 8.3, 1.2, 4.5, 5.5, CC_LIGHT)
        add_text_box(slide, info_box_text, 8.5, 1.4, 4.1, 5.1,
                     font_size=14, color=CC_BLACK, bold=True)

    if number:
        slide_number(slide, number)

    return slide


# ─── Layout 7: Discussion / Questions ────────────────────────────────────────
def discussion_slide(pptx, title, questions, number=None):
    """
    Heading + up to 6 numbered question items in 2 columns.
    Number badges: CC_PRIMARY filled rectangle + white number.
    questions: list of strings — max 6
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_PAGE)

    add_text_box(slide, title, 0.5, 0.3, 12.33, 0.7,
                 font_size=32, bold=True, color=CC_BLACK)

    qs     = questions[:6]
    half   = (len(qs) + 1) // 2
    left_qs  = qs[:half]
    right_qs = qs[half:]

    col_w  = 5.8
    col_x  = [0.5, 6.8]
    item_h = 1.4
    start_y = 1.3

    for col_i, col in enumerate([left_qs, right_qs]):
        for row_i, q in enumerate(col):
            y   = start_y + row_i * item_h
            x   = col_x[col_i]
            num = col_i * half + row_i + 1

            add_rect(slide, x, y, 0.45, 0.45, CC_PRIMARY)
            add_text_box(slide, str(num), x, y + 0.02, 0.45, 0.42,
                         font_size=14, bold=True, color=CC_WHITE,
                         align=PP_ALIGN.CENTER)
            add_text_box(slide, q, x + 0.6, y, col_w - 0.7, 1.2,
                         font_size=15, color=CC_BLACK)

    if number:
        slide_number(slide, number)

    return slide


# ─── Layout 8: Team / Facilitators ───────────────────────────────────────────
def team_grid(pptx, title, members, number=None):
    """
    Heading + 4 equal cards (CC_BASE), each with:
    - Circle avatar placeholder (CC_LIGHT fill, CC_PRIMARY border)
    - Initials inside circle
    - Name (bold CC_PRIMARY)
    - Title (black)
    members: list of dicts [{name, title?}, ...] — max 4
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_PAGE)

    add_text_box(slide, title, 0.5, 0.3, 12.33, 0.7,
                 font_size=32, bold=True, color=CC_BLACK)

    n       = min(len(members), 4)
    card_w  = 2.8
    gap     = 0.3
    total_w = n * card_w + (n - 1) * gap
    start_x = (13.33 - total_w) / 2
    card_top = 1.4
    card_h  = 5.2

    for i, m in enumerate(members[:n]):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, card_top, card_w, card_h, CC_BASE)

        # Circle avatar
        circle = slide.shapes.add_shape(
            9,  # oval
            Inches(x + (card_w - 1.4) / 2),
            Inches(card_top + 0.3),
            Inches(1.4),
            Inches(1.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = CC_LIGHT
        circle.line.color.rgb = CC_PRIMARY
        circle.line.width = Pt(2)

        initials = ""
        if m.get("name"):
            parts = m["name"].split()
            initials = "".join(p[0] for p in parts[:2]).upper()
        add_text_box(slide, initials,
                     x + (card_w - 1.4) / 2, card_top + 0.55,
                     1.4, 0.8,
                     font_size=24, bold=True, color=CC_BLACK,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, m.get("name", ""), x, card_top + 2.0,
                     card_w, 0.5,
                     font_size=16, bold=True, color=CC_BLACK,
                     align=PP_ALIGN.CENTER)

        if m.get("title"):
            add_text_box(slide, m.get("title", ""), x + 0.1, card_top + 2.5,
                         card_w - 0.2, 2.0,
                         font_size=12, color=CC_BLACK,
                         align=PP_ALIGN.CENTER)

    if number:
        slide_number(slide, number)

    return slide


# ─── Layout 9: Quote Slide ───────────────────────────────────────────────────
def quote_slide(pptx, quote, author="", number=None):
    """
    CC_BASE background, large decorative quote mark (CC_LIGHT),
    26pt CC_PRIMARY quote text, right-aligned author in CC_SECONDARY.
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_BASE)

    # Left accent bar
    add_rect(slide, 0.5, 1.5, 0.12, 4.5, CC_PRIMARY)

    # Open-quote decoration
    add_text_box(slide, '"', 0.8, 1.0, 1.5, 1.5,
                 font_size=120, bold=True, color=CC_LIGHT)

    # Quote text
    add_text_box(slide, quote, 1.2, 2.0, 9.0, 3.5,
                 font_size=26, bold=False, color=CC_BLACK,
                 font_name=FONT_BODY)

    if author:
        add_text_box(slide, f"— {author}", 1.2, 5.6, 9.0, 0.5,
                     font_size=16, bold=True, color=CC_BLACK,
                     align=PP_ALIGN.RIGHT)

    if number:
        slide_number(slide, number)

    return slide


# ─── Layout 10: Closing Slide ─────────────────────────────────────────────────
def closing_slide(pptx, text="THANK YOU", subtitle="", number=None):
    """
    Dark closing: CC_PRIMARY background, 54pt white bold text,
    large decorative star watermark, optional subtitle.
    """
    slide = _add_slide(pptx)
    set_slide_bg(slide, CC_PRIMARY)

    # Star watermark
    add_text_box(slide, "★", 5.5, 0.5, 2.3, 2.5,
                 font_size=150, bold=True,
                 color=RGBColor(0x7A, 0x25, 0x1F),
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, text, 1.0, 2.8, 11.33, 1.5,
                 font_size=54, bold=True, color=CC_WHITE,
                 align=PP_ALIGN.CENTER)

    if subtitle:
        add_text_box(slide, subtitle, 1.5, 4.5, 10.33, 0.8,
                     font_size=24, color=CC_LIGHT,
                     align=PP_ALIGN.CENTER)

    if number:
        slide_number(slide, number, dark_bg=True)

    return slide
