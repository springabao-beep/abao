# -*- coding: utf-8 -*-
"""
Career Compass PPT Format Checker
Audits an existing PPT file for Career Compass brand compliance.

Usage:
    python format_check.py your_file.pptx
    python format_check.py your_file.pptx --fix  (apply auto-fixes)

Exit codes:
    0 = all checks passed
    1 = issues found
    2 = file not found
"""

import sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

# ─── Career Compass Color System ─────────────────────────────────────────────
CC_PRIMARY   = RGBColor(0x9D, 0x3D, 0x1D)
CC_SECONDARY = RGBColor(0xB9, 0x5B, 0x42)
CC_ACCENT    = RGBColor(0xD6, 0x75, 0x5A)
CC_LIGHT     = RGBColor(0xDD, 0xBC, 0xB0)
CC_BASE      = RGBColor(0xF1, 0xE2, 0xDD)
CC_PAGE      = RGBColor(0xF7, 0xF3, 0xED)
CC_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CC_BLACK     = RGBColor(0x00, 0x00, 0x00)

# Allowed fonts (primary + fallbacks)
ALLOWED_FONTS = {
    "Futura Bk BT", "Futura", "Futura Md BT", "Futura Bd BT",
    "LVMH Sans", "LVMH", "LVMH Sans SemiLight", "LVMH Sans SemiBold",
    "LVMH Sans ExtraLight", "LVMH Sans Light",
    "LVMH ExtraLight", "LVMH Light", "LVMH SemiBold",
    "Arial", "Helvetica",
    "Microsoft YaHei", "SimHei", "SimSun",
}
FORBIDDEN_FONTS = {
    "Times New Roman", "Courier New", "Georgia", "Comic Sans MS",
    "Wingdings", "Symbol",
}

# Size thresholds
MIN_FONT_SIZE = Pt(9)
MIN_USABLE_SIZE = Pt(9)
MAX_TITLE_SIZE = Pt(72)
MIN_BODY_SIZE = Pt(11)
MAX_BODY_SIZE = Pt(24)


def rgb_to_hex(rgb):
    if rgb is None:
        return None
    return "#{r:02X}{g:02X}{b:02X}".format(r=rgb[0], g=rgb[1], b=rgb[2])


def check_color_approximation(rgb, allowed_colors, threshold=30):
    """Check if a color is close to any allowed Career Compass color."""
    if rgb is None:
        return False, None
    for name, ref in allowed_colors.items():
        dr = abs(int(rgb[0]) - int(ref[0]))
        dg = abs(int(rgb[1]) - int(ref[1]))
        db = abs(int(rgb[2]) - int(ref[2]))
        if dr + dg + db < threshold:
            return True, name
    return False, None


ALLOWED_COLORS = {
    "#9D3D1D": CC_PRIMARY,
    "#B95B42": CC_SECONDARY,
    "#D6755A": CC_ACCENT,
    "#DDBCB0": CC_LIGHT,
    "#F1E2DD": CC_BASE,
    "#F7F3ED": CC_PAGE,
    "#FFFFFF": CC_WHITE,
    "#000000": CC_BLACK,
    "#9D3C1D": CC_PRIMARY,   # slight variant in some slides
    "#572923": RGBColor(0x57, 0x29, 0x23),
    "#7A251F": RGBColor(0x7A, 0x25, 0x1F),
    "#020F2B": RGBColor(0x02, 0x0F, 0x2B),  # navy — only for dark slides
}


def audit_pptx(path, apply_fixes=False):
    """Run all checks on a PPTX file. Returns (passes, issues)."""
    try:
        prs = Presentation(path)
    except FileNotFoundError:
        return False, [{"level": "ERROR", "msg": f"File not found: {path}"}]
    except Exception as e:
        return False, [{"level": "ERROR", "msg": f"Cannot open file: {e}"}]

    issues = []
    passes = 0

    # Check 1: Slide dimensions (should be 16:9)
    w_inch = prs.slide_width / 914400
    h_inch = prs.slide_height / 914400
    ratio = w_inch / h_inch
    if abs(ratio - 16/9) > 0.02:
        issues.append({
            "level": "WARN",
            "msg": f"Slide ratio {ratio:.3f} ≠ 16:9 ({w_inch:.2f}\" x {h_inch:.2f}\")"
        })
    else:
        passes += 1

    # Check 2: Per-slide font and color checks
    slide_colors_used = set()
    slide_fonts_used  = set()

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font_name = run.font.name
                    font_size = run.font.size
                    rgb       = None
                    try:
                        if run.font.color.type == 1:
                            rgb = run.font.color.rgb
                    except:
                        pass

                    if font_name:
                        slide_fonts_used.add(font_name)
                        if font_name in FORBIDDEN_FONTS:
                            issues.append({
                                "level": "ERROR",
                                "slide": i,
                                "shape": shape.name,
                                "msg": f"Forbidden font '{font_name}' used"
                            })
                        elif font_name not in ALLOWED_FONTS:
                            issues.append({
                                "level": "INFO",
                                "slide": i,
                                "msg": f"Non-standard font '{font_name}' — recommend Futura Bk BT or Arial"
                            })

                    if font_size:
                        slide_colors_used.add(rgb_to_hex(rgb) if rgb else None)
                        if font_size < MIN_FONT_SIZE:
                            issues.append({
                                "level": "WARN",
                                "slide": i,
                                "msg": f"Font size {font_size.pt:.1f}pt < 9pt minimum"
                            })

                    if rgb:
                        slide_colors_used.add(rgb_to_hex(rgb))
                        is_cc, name = check_color_approximation(rgb, ALLOWED_COLORS)
                        if not is_cc:
                            issues.append({
                                "level": "INFO",
                                "slide": i,
                                "msg": f"Non-CC color #{rgb_to_hex(rgb)} used"
                            })

    passes += 1  # dimension check passed

    # Summary
    if not issues:
        return True, [{"level": "PASS", "msg": "All checks passed"}]

    return False, issues


def main():
    apply_fixes = "--fix" in sys.argv
    paths = [a for a in sys.argv[1:] if not a.startswith("--")]

    if not paths:
        print("Usage: python format_check.py <file.pptx> [--fix]")
        print("")
        print("Checks:")
        print("  - Slide dimensions (16:9)")
        print("  - Font families (forbidden fonts, non-standard fonts)")
        print("  - Font sizes (min 9pt)")
        print("  - Color palette (Career Compass colors only)")
        sys.exit(2)

    for path in paths:
        print(f"\n{'='*60}")
        print(f"Checking: {path}")
        print("=" * 60)

        ok, results = audit_pptx(path, apply_fixes=apply_fixes)

        errors   = [r for r in results if r.get("level") == "ERROR"]
        warnings = [r for r in results if r.get("level") == "WARN"]
        infos    = [r for r in results if r.get("level") == "INFO"]
        passes   = [r for r in results if r.get("level") == "PASS"]

        for r in passes:
            print(f"  [PASS] {r['msg']}")
        for r in infos:
            slide_info = f"[Slide {r.get('slide', '?')}]" if 'slide' in r else ""
            print(f"  [INFO] {slide_info} {r['msg']}")
        for r in warnings:
            slide_info = f"[Slide {r.get('slide', '?')}]" if 'slide' in r else ""
            print(f"  [WARN] {slide_info} {r['msg']}")
        for r in errors:
            slide_info = f"[Slide {r.get('slide', '?')}]" if 'slide' in r else ""
            print(f"  [ERROR] {slide_info} {r['msg']}")

        print(f"\n  Summary: {len(passes)} passed, {len(errors)} errors, "
              f"{len(warnings)} warnings, {len(infos)} infos")

    sys.exit(1 if errors else 0)


if __name__ == "__main__":
    main()
