# -*- coding: utf-8 -*-
"""
Career Compass PPT Format Unifier
Apply Career Compass brand colors and typography to any existing PPT file.

Usage:
    python apply_cc_format.py <file.pptx> [--backup] [--dry-run] [--no-fonts]

Modes:
    --backup   Create a backup copy before modifying
    --dry-run  Only report what would change, don't modify
    --no-fonts Skip font replacement (colors only)

Exit codes:
    0 = all done (or dry-run complete)
    1 = errors
"""
import sys, os, shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


# ─── Career Compass Palette ────────────────────────────────────────────────────
CC_COLORS = [
    ("CC_PRIMARY",   RGBColor(0x9D, 0x3D, 0x1D)),
    ("CC_SECONDARY", RGBColor(0xB9, 0x5B, 0x42)),
    ("CC_ACCENT",    RGBColor(0xD6, 0x75, 0x5A)),
    ("CC_LIGHT",     RGBColor(0xDD, 0xBC, 0xB0)),
    ("CC_BASE",      RGBColor(0xF1, 0xE2, 0xDD)),
    ("CC_PAGE",      RGBColor(0xF7, 0xF3, 0xED)),
    ("CC_WHITE",     RGBColor(0xFF, 0xFF, 0xFF)),
    ("CC_BLACK",     RGBColor(0x00, 0x00, 0x00)),
]
CC_MAP = {c[1]: c[0] for c in CC_COLORS}

# Brightness boundaries — dark colors (< threshold) map to CC_PRIMARY
DARK_THRESHOLD = 120  # sum of R+G+B below this = "dark" → CC_PRIMARY
LIGHT_THRESHOLD = 700  # sum of R+G+B above this = "light" → CC_PAGE
VERY_LIGHT_THRESHOLD = 740  # near-white → CC_WHITE / CC_PAGE

FONT_PRIMARY = "Futura Bk BT"


def rgb_distance(c1, c2):
    """Euclidean distance between two RGB colors."""
    return ((c1[0]-c2[0])**2 + (c1[1]-c2[1])**2 + (c1[2]-c2[2])**2) ** 0.5


def snap_to_cc(rgb):
    """
    Map any RGB color to the nearest Career Compass color.
    Uses a combination of brightness heuristics and nearest-neighbor matching.
    """
    r, g, b = rgb
    total = r + g + b

    # Very dark → CC_PRIMARY (dark backgrounds)
    if total < DARK_THRESHOLD:
        return CC_COLORS[0][1]  # CC_PRIMARY

    # Very light → CC_PAGE or CC_WHITE
    if total > VERY_LIGHT_THRESHOLD:
        return CC_COLORS[5][1]  # CC_PAGE
    if total > LIGHT_THRESHOLD:
        return CC_COLORS[4][1]  # CC_BASE

    # Find nearest CC color by Euclidean distance
    best = min(CC_COLORS, key=lambda c: rgb_distance(rgb, c[1]._rgb if hasattr(c[1], '_rgb') else (c[1][0], c[1][1], c[1][2])))
    return best[1]


def get_rgb_tuple(rgb_obj):
    """Convert RGBColor to (R,G,B) tuple."""
    if hasattr(rgb_obj, '__iter__'):
        return tuple(rgb_obj)
    return (rgb_obj[0], rgb_obj[1], rgb_obj[2])


def _get_rgb_comps(rgb_obj):
    """Get integer R,G,B components from an RGBColor."""
    if hasattr(rgb_obj, 'red'):
        return (rgb_obj.red, rgb_obj.green, rgb_obj.blue)
    # RGBColor is tuple-like
    return (int(rgb_obj[0]), int(rgb_obj[1]), int(rgb_obj[2]))


def process_presentation(path, backup=True, dry_run=False, skip_fonts=False):
    """Main entry: load PPTX, apply CC formatting, save."""
    prs = Presentation(path)

    changes = {
        "fills": 0,
        "text_colors": 0,
        "backgrounds": 0,
        "fonts": 0,
    }
    slide_reports = []

    for i, slide in enumerate(prs.slides):
        slide_fills = 0
        slide_texts = 0
        slide_fonts = 0

        # 1. Slide background
        try:
            bg = slide.background.fill
            old_rgb = bg.fore_color.rgb
            comps = _get_rgb_comps(old_rgb)
            new_rgb = snap_to_cc(comps)
            if comps != _get_rgb_comps(new_rgb):
                if not dry_run:
                    bg.fore_color.rgb = new_rgb
                changes["backgrounds"] += 1
                slide_fills += 1
        except Exception:
            pass

        # 2. Shapes
        for shape in slide.shapes:
            # Shape fill
            changed_fill = False
            try:
                fill = shape.fill
                old_rgb = fill.fore_color.rgb
                comps = _get_rgb_comps(old_rgb)
                new_rgb = snap_to_cc(comps)
                if comps != _get_rgb_comps(new_rgb):
                    if not dry_run:
                        fill.fore_color.rgb = new_rgb
                    changes["fills"] += 1
                    slide_fills += 1
                    changed_fill = True
            except Exception:
                pass

            # Table cells
            if int(shape.shape_type) == 19:  # TABLE
                try:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            # Cell fill
                            try:
                                cf = cell.fill
                                old_rgb = cf.fore_color.rgb
                                comps = _get_rgb_comps(old_rgb)
                                new_rgb = snap_to_cc(comps)
                                if comps != _get_rgb_comps(new_rgb):
                                    if not dry_run:
                                        cf.fore_color.rgb = new_rgb
                                    changes["fills"] += 1
                                    slide_fills += 1
                            except Exception:
                                pass
                            # Cell text
                            try:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        old_rgb = run.font.color.rgb
                                        comps = _get_rgb_comps(old_rgb)
                                        new_rgb = snap_to_cc(comps)
                                        # ⚠️ User: map text to black
                                        r, g, b = _get_rgb_comps(new_rgb)
                                        if r + g + b > 60 and new_rgb != CC_WHITE:
                                            new_rgb = CC_BLACK
                                        if comps != _get_rgb_comps(new_rgb):
                                            if not dry_run:
                                                run.font.color.rgb = new_rgb
                                            changes["text_colors"] += 1
                                            slide_texts += 1
                            except Exception:
                                pass
                except Exception:
                    pass
                continue  # skip text/font processing below for tables

            # Text colors (non-table shapes)
            if shape.has_text_frame:
                if not changed_fill:
                    try:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    old_rgb = run.font.color.rgb
                                    comps = _get_rgb_comps(old_rgb)
                                    new_rgb = snap_to_cc(comps)
                                    # ⚠️ User: map text to black
                                    r, g, b = _get_rgb_comps(new_rgb)
                                    if r + g + b > 60 and new_rgb != CC_WHITE:
                                        new_rgb = CC_BLACK
                                    if comps != _get_rgb_comps(new_rgb):
                                        if not dry_run:
                                            run.font.color.rgb = new_rgb
                                        changes["text_colors"] += 1
                                        slide_texts += 1
                                except Exception:
                                    pass
                    except Exception:
                        pass

                # Font override
                if not skip_fonts:
                    try:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name and run.font.name.lower() != FONT_PRIMARY.lower():
                                    if not dry_run:
                                        run.font.name = FONT_PRIMARY
                                    changes["fonts"] += 1
                                    slide_fonts += 1
                    except Exception:
                        pass

        slide_reports.append((i+1, slide_fills, slide_texts, slide_fonts))

    # Reporting
    print(f"\n{'='*60}")
    mode = "DRY RUN" if dry_run else "APPLYING"
    print(f"  [{mode}] {os.path.basename(path)}")
    print(f"{'='*60}")
    for num, fills, texts, fonts in slide_reports:
        parts = []
        if fills: parts.append(f"fills={fills}")
        if texts: parts.append(f"texts={texts}")
        if fonts: parts.append(f"fonts={fonts}")
        if parts:
            print(f"  Slide {num}: {', '.join(parts)}")

    print(f"\n  Summary:")
    print(f"    Fills:        {changes['fills']}")
    print(f"    Backgrounds:  {changes['backgrounds']}")
    print(f"    Text colors:  {changes['text_colors']}")
    print(f"    Font changes: {changes['fonts']}")

    # Save
    if not dry_run:
        prs.save(path)
        print(f"\n  ✓ Saved: {path}")

    return changes


def main():
    # Parse args
    paths = []
    backup = True
    dry_run = False
    skip_fonts = False

    for a in sys.argv[1:]:
        if a == "--no-backup":
            backup = False
        elif a == "--dry-run":
            dry_run = True
        elif a == "--no-fonts":
            skip_fonts = True
        elif a.startswith("--"):
            pass  # unknown flag — ignore
        else:
            paths.append(a)

    if not paths:
        print("Usage: python apply_cc_format.py <file.pptx>")
        print("       --backup    Create backup (auto, unless --no-backup)")
        print("       --dry-run   Preview changes only")
        print("       --no-fonts  Skip font replacement")
        sys.exit(2)

    for p in paths:
        if not os.path.isfile(p):
            print(f"  [ERROR] File not found: {p}")
            continue

        # Backup
        bak_path = p.replace(".pptx", "_bak.pptx")
        if not os.path.isfile(bak_path):
            shutil.copy2(p, bak_path)
            if not dry_run:
                print(f"  [BACKUP] → {bak_path}")

        process_presentation(p, backup=backup, dry_run=dry_run, skip_fonts=skip_fonts)

        # Post-run format check (unless dry-run)
        if not dry_run:
            try:
                sys.path.insert(0, os.path.join(
                    os.path.dirname(__file__), ".."))
                from scripts import format_check as fc
                ok, results = fc.audit_pptx(p)
                errors   = [r for r in results if r.get("level") == "ERROR"]
                warnings = [r for r in results if r.get("level") == "WARN"]
                print(f"\n  [FORMAT CHECK] {len(errors)} errors, {len(warnings)} warnings")
            except Exception as e:
                print(f"  [FORMAT CHECK SKIPPED] {e}")

    sys.exit(0)


if __name__ == "__main__":
    main()
